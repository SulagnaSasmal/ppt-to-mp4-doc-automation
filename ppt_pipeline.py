import win32com.client
import subprocess
import os
import time
import shutil
from pathlib import Path
from functools import lru_cache
import requests
from typing import Callable, Dict, Any, Optional

# Use numeric constant to avoid makepy dependency when constants are unavailable
PPT_ADVANCE_USE_TIMINGS = 3  # ppSlideShowUseSlideTimings
DEFAULT_SILENCE_HEAD = 0.0
DEFAULT_SILENCE_TAIL = 0.0
VOICE = "en-US-JennyNeural"
FFMPEG_BIN_DIR_ENV = "FFMPEG_BIN_DIR"
DEFAULT_PIPELINE_SETTINGS = {
    "voice": VOICE,
    "speaking_rate": "0%",
    "resolution": 1080,
    "fps": 30,
    "quality": 100,
}


@lru_cache(maxsize=2)
def resolve_media_tool(tool_name: str) -> str:
    """Resolve ffmpeg/ffprobe executable from PATH, env override, or winget folders."""
    from_path = shutil.which(tool_name)
    if from_path:
        return from_path

    candidates = []

    custom_bin_dir = os.environ.get(FFMPEG_BIN_DIR_ENV)
    if custom_bin_dir:
        candidates.append(Path(custom_bin_dir) / f"{tool_name}.exe")

    local_appdata = os.environ.get("LOCALAPPDATA")
    if local_appdata:
        winget_packages = Path(local_appdata) / "Microsoft" / "WinGet" / "Packages"
        if winget_packages.exists():
            candidates.extend(winget_packages.glob(f"**/{tool_name}.exe"))

    for candidate in candidates:
        if candidate.exists():
            return str(candidate)

    raise FileNotFoundError(
        f"{tool_name} executable not found. Install FFmpeg and ensure '{tool_name}' is on PATH, "
        f"or set {FFMPEG_BIN_DIR_ENV} to the FFmpeg bin directory."
    )


def probe_duration_seconds(path: Path) -> float:
    """Return media duration in seconds using ffprobe."""
    result = subprocess.check_output([
        resolve_media_tool("ffprobe"), "-v", "error",
        "-show_entries", "format=duration",
        "-of", "default=nk=1:nw=1",
        str(path)
    ], text=True).strip()
    return float(result)


def wait_for_readable(path: Path, attempts: int = 30, delay: float = 1.0) -> None:
    """Wait until a file can be opened for reading (handles lingering locks)."""
    for _ in range(attempts):
        try:
            with open(path, "rb"):
                return
        except OSError:
            time.sleep(delay)
    raise FileNotFoundError(f"File not readable (locked?): {path}")


def safe_unlink(path: Path, attempts: int = 5, delay: float = 0.5) -> None:
    """Delete a file with retries to handle transient locks."""
    for i in range(attempts):
        if not path.exists():
            return
        try:
            path.unlink()
            return
        except OSError:
            time.sleep(delay)
    raise PermissionError(f"Could not delete locked file: {path}")


def safe_close_presentation(presentation) -> None:
    """Best-effort close for PowerPoint presentation objects."""
    if not presentation:
        return
    try:
        presentation.Close()
        return
    except Exception:
        pass

    try:
        presentation.Saved = True
        presentation.Close()
    except Exception as e:
        print(f"[WARNING] Failed to close presentation cleanly: {e}")


def safe_quit_powerpoint(ppt) -> None:
    """Best-effort quit for PowerPoint application object."""
    if not ppt:
        return
    try:
        ppt.Quit()
    except Exception as e:
        print(f"[WARNING] Failed to quit PowerPoint cleanly: {e}")


def normalize_pipeline_settings(settings: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    merged = dict(DEFAULT_PIPELINE_SETTINGS)
    if settings:
        merged.update({k: v for k, v in settings.items() if v is not None})

    merged["voice"] = str(merged["voice"]).strip() or VOICE
    merged["speaking_rate"] = str(merged["speaking_rate"]).strip() or "0%"
    merged["resolution"] = max(240, min(2160, int(merged["resolution"])))
    merged["fps"] = max(1, min(60, int(merged["fps"])))
    merged["quality"] = max(1, min(100, int(merged["quality"])))
    return merged


def extract_slide_notes(ppt_path: str) -> list[dict]:
    from pptx import Presentation

    presentation = Presentation(str(Path(ppt_path).resolve()))
    notes = []
    for index, slide in enumerate(presentation.slides, start=1):
        text = ""
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            text = (slide.notes_slide.notes_text_frame.text or "").strip()
        notes.append({
            "slide": index,
            "text": text,
            "has_notes": bool(text),
        })
    return notes


def run_pipeline(
    ppt_path: str,
    out_dir: str,
    settings: Optional[Dict[str, Any]] = None,
    progress_cb: Optional[Callable[[str, int, str], None]] = None,
) -> Dict[str, Any]:
    """
    Runs:
    - PPT → animated MP4
    - Notes → Azure TTS
    - FFmpeg mux
    Returns final MP4 path
    """
    ppt_path_obj = Path(ppt_path).resolve()
    out_dir_path = Path(out_dir).resolve()
    out_dir_path.mkdir(parents=True, exist_ok=True)
    pipeline_settings = normalize_pipeline_settings(settings)

    def report(stage: str, progress: int, message: str) -> None:
        if progress_cb:
            progress_cb(stage, progress, message)

    total_start = time.perf_counter()

    video_raw = out_dir_path / "video_raw.mp4"
    narration = out_dir_path / "narration.mp3"
    final_video = out_dir_path / "final.mp4"

    # Check for required Azure credentials
    azure_key = os.environ.get("AZURE_TTS_KEY")
    azure_region = os.environ.get("AZURE_TTS_REGION")
    
    if not azure_key or not azure_region:
        raise ValueError(
            "Missing Azure TTS credentials. Please set environment variables:\n"
            "  AZURE_TTS_KEY=<your-key>\n"
            "  AZURE_TTS_REGION=<region> (e.g., 'eastus')\n"
            "Before starting the server."
        )
    
    tts_url = f"https://{azure_region}.tts.speech.microsoft.com/cognitiveservices/v1"
    tts_headers = {
        "Ocp-Apim-Subscription-Key": azure_key,
        "Content-Type": "application/ssml+xml",
        "X-Microsoft-OutputFormat": "audio-24khz-160kbitrate-mono-mp3",
    }

    report("notes", 12, "Extracting slide notes")
    extracted_notes = extract_slide_notes(str(ppt_path_obj))
    slide_notes = [(item["slide"], item["text"]) for item in extracted_notes]

    # 2️⃣ Azure TTS per slide
    report("tts", 35, "Generating Azure TTS from notes")
    tts_start = time.perf_counter()
    audio_files = []
    audio_durations = {}
    for i, text in slide_notes:
        if not text:
            continue

        ssml = f"""
<speak version='1.0' xml:lang='en-US'>
  <voice name='{pipeline_settings['voice']}'><prosody rate='{pipeline_settings['speaking_rate']}'>{text}</prosody></voice>
</speak>
"""

        out_mp3 = (out_dir_path / f"slide_{i:02}.mp3").resolve()
        safe_unlink(out_mp3)
        r = requests.post(tts_url, headers=tts_headers, data=ssml.encode("utf-8"))
        r.raise_for_status()
        out_mp3.write_bytes(r.content)
        audio_files.append(out_mp3)
        audio_durations[i] = probe_duration_seconds(out_mp3)

    if not audio_files:
        raise ValueError("No speaker notes found for TTS. Add notes to at least one slide before converting.")

    tts_seconds = round(time.perf_counter() - tts_start, 2)

    # 3️⃣ Open PPT and set slide timings to match narration durations (plus padding)
    report("ppt_export", 60, "Exporting animated video via PowerPoint")
    ppt_export_start = time.perf_counter()
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    ppt.Visible = True

    presentation = ppt.Presentations.Open(
        str(ppt_path_obj),
        WithWindow=True
    )

    for slide_index, slide in enumerate(presentation.Slides, start=1):
        slide.SlideShowTransition.AdvanceOnTime = True
        slide.SlideShowTransition.AdvanceOnClick = False
        slide.SlideShowTransition.Duration = 0.0
        duration = audio_durations.get(slide_index, 2.0) + DEFAULT_SILENCE_HEAD + DEFAULT_SILENCE_TAIL
        slide.SlideShowTransition.AdvanceTime = duration

    presentation.SlideShowSettings.AdvanceMode = PPT_ADVANCE_USE_TIMINGS
    presentation.SlideShowSettings.ShowWithAnimation = True
    presentation.SlideShowSettings.ShowWithNarration = True

    # Save a clean copy to avoid compatibility/protection quirks, then reopen it
    tmp_pptx = (out_dir_path / f"tmp_for_video_{int(time.time())}.pptx").resolve()
    safe_unlink(tmp_pptx)
    try:
        presentation.SaveCopyAs(str(tmp_pptx))
    except Exception as e:
        safe_close_presentation(presentation)
        safe_quit_powerpoint(ppt)
        raise RuntimeError(f"PowerPoint could not save a temp copy: {e}") from e

    safe_close_presentation(presentation)
    presentation = ppt.Presentations.Open(str(tmp_pptx), WithWindow=True)

    # Re-apply timings on reopened copy
    for slide_index, slide in enumerate(presentation.Slides, start=1):
        slide.SlideShowTransition.AdvanceOnTime = True
        slide.SlideShowTransition.AdvanceOnClick = False
        slide.SlideShowTransition.Duration = 0.0
        duration = audio_durations.get(slide_index, 2.0) + DEFAULT_SILENCE_HEAD + DEFAULT_SILENCE_TAIL
        slide.SlideShowTransition.AdvanceTime = duration

    presentation.SlideShowSettings.AdvanceMode = PPT_ADVANCE_USE_TIMINGS
    presentation.SlideShowSettings.ShowWithAnimation = True
    presentation.SlideShowSettings.ShowWithNarration = True

    # Persist timings before rendering
    presentation.Save()

    safe_unlink(video_raw)
    safe_unlink(narration)
    safe_unlink(final_video)

    # 4️⃣ Export video with slide-based timings
    presentation.CreateVideo(
        str(video_raw),
        UseTimingsAndNarrations=True,
        DefaultSlideDuration=1,
        VertResolution=pipeline_settings["resolution"],
        FramesPerSecond=pipeline_settings["fps"],
        Quality=pipeline_settings["quality"]
    )

    status = presentation.CreateVideoStatus
    while status == 1:  # In progress
        time.sleep(2)
        status = presentation.CreateVideoStatus

    if status != 3:  # 3 = Success
        # Fallback to SaveAs MP4 (ppSaveAsMP4 = 39). Some builds reject CreateVideo but accept SaveAs.
        try:
            presentation.SaveAs(str(video_raw), FileFormat=39)
            for _ in range(30):  # wait up to ~60s
                if video_raw.exists():
                    status = 3
                    break
                time.sleep(2)
        except Exception as e:
            safe_close_presentation(presentation)
            safe_quit_powerpoint(ppt)
            raise RuntimeError(
                "PowerPoint failed to render video (status={}). SaveAs fallback also failed: {}".format(status, e)
            ) from e

    if status != 3:
        safe_close_presentation(presentation)
        safe_quit_powerpoint(ppt)
        raise RuntimeError(
            "PowerPoint failed to render video (status={}). "
            "Tried CreateVideo and SaveAs MP4. Check for hidden dialogs or re-run once manually from UI to unblock.".format(status)
        )

    if not video_raw.exists():
        raise FileNotFoundError(f"PowerPoint did not produce the video file: {video_raw}")
    print("[OK] Animated video created:", video_raw)

    # Wait until the MP4 is actually written and non-empty
    for _ in range(60):  # up to ~120s
        try:
            size = video_raw.stat().st_size
        except FileNotFoundError:
            size = 0
        if size > 0:
            break
        time.sleep(2)
    else:
        try:
            # Fallback: try SaveAs MP4 if CreateVideo wrote a zero-byte file
            presentation.SaveAs(str(video_raw), FileFormat=39)
            for _ in range(30):
                try:
                    size = video_raw.stat().st_size
                except FileNotFoundError:
                    size = 0
                if size > 0:
                    break
                time.sleep(2)
        except Exception:
            safe_close_presentation(presentation)
            safe_quit_powerpoint(ppt)
            raise FileNotFoundError(f"PowerPoint produced an empty/locked file and SaveAs retry failed: {video_raw}")

        if size <= 0:
            safe_close_presentation(presentation)
            safe_quit_powerpoint(ppt)
            raise FileNotFoundError(f"PowerPoint produced an empty/locked file: {video_raw}")

    print(f"[OK] Video file size: {video_raw.stat().st_size} bytes")

    # Ensure MP4 is readable (release any lingering lock)
    wait_for_readable(video_raw, attempts=30, delay=1.0)

    # Log intended total timed duration vs audio sum for sanity
    total_slide_seconds = sum(audio_durations.get(i, 2.0) + DEFAULT_SILENCE_HEAD + DEFAULT_SILENCE_TAIL for i, _ in slide_notes)
    try:
        video_seconds = probe_duration_seconds(video_raw)
        print(f"[INFO] Intended slide duration sum: {total_slide_seconds:.2f}s; Rendered video: {video_seconds:.2f}s")
    except Exception:
        pass

    # Close PowerPoint with retry and error handling
    safe_close_presentation(presentation)
    safe_quit_powerpoint(ppt)
    ppt_export_seconds = round(time.perf_counter() - ppt_export_start, 2)

    # 5️⃣ Concatenate audio clips
    report("mux", 85, "FFmpeg muxing")
    mux_start = time.perf_counter()
    audio_list = (out_dir_path / "audio_list.txt").resolve()
    with audio_list.open("w", encoding="utf-8") as f:
        for a in audio_files:
            f.write(f"file '{a.as_posix()}'\n")

    ffmpeg_exe = resolve_media_tool("ffmpeg")

    subprocess.run([
        ffmpeg_exe, "-y",
        "-f", "concat",
        "-safe", "0",
        "-i", str(audio_list),
        "-c", "copy",
        str(narration)
    ], check=True)

    # Ensure narration file is present before mux
    if not narration.exists():
        raise FileNotFoundError(f"Narration not found: {narration}")

    # 6️⃣ FFmpeg mux (KEEP ANIMATIONS)
    subprocess.run([
        ffmpeg_exe, "-y",
        "-i", str(video_raw),
        "-i", str(narration),
        "-map", "0:v:0",
        "-map", "1:a:0",
        "-c:v", "copy",
        "-c:a", "aac",
        "-shortest",
        str(final_video)
    ], check=True)

    mux_seconds = round(time.perf_counter() - mux_start, 2)
    total_seconds = round(time.perf_counter() - total_start, 2)
    telemetry = {
        "tts_seconds": tts_seconds,
        "ppt_export_seconds": ppt_export_seconds,
        "mux_seconds": mux_seconds,
        "total_seconds": total_seconds,
    }

    print("[OK] DONE ->", final_video)
    return {
        "final_video": str(final_video),
        "telemetry": telemetry,
        "settings": pipeline_settings,
        "notes_with_text": len([n for _, n in slide_notes if n]),
        "slides_total": len(slide_notes),
    }
